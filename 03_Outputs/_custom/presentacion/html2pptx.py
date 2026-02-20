#!/usr/bin/env python3
"""
html2pptx.py — Converts TIMining CORE HTML presentation to editable PPTX.

Parses index.html (34 slides), generates a clean light-mode PowerPoint
preserving layout structure and text hierarchy. Editable, not pixel-perfect.

Usage: python3 html2pptx.py
Output: TIMining_CORE_Presentacion.pptx
"""

import re
import os
from pathlib import Path
from bs4 import BeautifulSoup, Tag
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Light-mode palette
# ---------------------------------------------------------------------------
BG = RGBColor(0xF8, 0xF9, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_BODY = RGBColor(0x33, 0x33, 0x33)
TEXT_MUTED = RGBColor(0x88, 0x88, 0x88)
ACCENT = RGBColor(0x00, 0x88, 0xAA)  # teal
COLOR_D = RGBColor(0xD3, 0x2F, 0x2F)
COLOR_A = RGBColor(0xE6, 0x51, 0x00)
COLOR_R = ACCENT
CARD_BG = WHITE
CARD_BORDER = RGBColor(0xE0, 0xE0, 0xE0)
IDEMAX_RED = RGBColor(0xF9, 0x35, 0x45)

# ---------------------------------------------------------------------------
# Dimensions (16:9)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
ML = Inches(0.7)  # margin left
MR = Inches(0.7)
MT = Inches(0.15)
CW = SLIDE_W - ML - MR  # content width

BASE_DIR = Path(__file__).parent
HTML_FILE = BASE_DIR / "index.html"
OUTPUT_FILE = BASE_DIR / "TIMining_CORE_Presentacion_v2.pptx"

TOTAL_SLIDES = 34


# ===========================================================================
# Helpers
# ===========================================================================

def clean(el):
    if el is None:
        return ""
    return re.sub(r'\s+', ' ', el.get_text(separator=" ", strip=True)).strip()


def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def tb(slide, left, top, w, h):
    return slide.shapes.add_textbox(left, top, w, h)


def run(p, text, sz=12, color=TEXT_BODY, bold=False, italic=False, font="Calibri"):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return r


def header(slide, num):
    """Compact header: brand left, number right."""
    t = tb(slide, ML, MT, Inches(2), Inches(0.35))
    p = t.text_frame.paragraphs[0]
    run(p, "TIMining ", sz=11, color=TEXT_MUTED, font="Consolas")
    run(p, "CORE", sz=11, color=ACCENT, bold=True, font="Consolas")

    t2 = tb(slide, SLIDE_W - MR - Inches(1), MT, Inches(1), Inches(0.35))
    p2 = t2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run(p2, f"{num:02d} / {TOTAL_SLIDES}", sz=10, color=TEXT_MUTED, font="Consolas")


def title(slide, text, top=Inches(0.5), sz=20, color=TEXT_DARK):
    t = tb(slide, ML, top, CW, Inches(0.45))
    t.text_frame.word_wrap = True
    p = t.text_frame.paragraphs[0]
    run(p, text, sz=sz, color=color, bold=True)


def subtitle(slide, text, top=Inches(0.95), sz=12, color=TEXT_MUTED):
    if not text:
        return
    t = tb(slide, ML, top, CW, Inches(0.4))
    t.text_frame.word_wrap = True
    p = t.text_frame.paragraphs[0]
    run(p, text, sz=sz, color=color)


def add_image_safe(slide, img_path, left, top, width=None, height=None):
    """Try to add an image, return the shape if successful, None otherwise."""
    path = BASE_DIR / img_path if not os.path.isabs(str(img_path)) else Path(img_path)
    if not path.exists():
        return None
    try:
        kwargs = {"image_file": str(path), "left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        return slide.shapes.add_picture(**kwargs)
    except Exception:
        return None


def extract_body(el):
    """Extract body text from a card element, handling <p>, <ul>/<li>,
    <strong>/<span>, and recap-section structures."""
    parts = []
    footer = ""

    # Recap sections (STOP/START blocks with <ul>/<li>)
    recap_sections = el.find_all("div", class_="recap-section")
    if recap_sections:
        for sec in recap_sections:
            label_el = sec.find("div", class_="recap-section-label")
            label = clean(label_el) if label_el else ""
            items = [clean(li) for li in sec.find_all("li")]
            if label and items:
                parts.append(f"{label}: " + " | ".join(items))
            elif items:
                parts.extend(["- " + item for item in items])
        return "\n".join(parts), footer

    # Standard: collect <p> and <ul> in document order
    for child in el.children:
        if not isinstance(child, Tag):
            continue
        if child.name == "p":
            txt = clean(child)
            if not txt:
                continue
            style = child.get("style", "")
            if "font-size: 11px" in style or "font-size: 10px" in style:
                footer = txt
            else:
                parts.append(txt)
        elif child.name == "ul":
            items = [clean(li) for li in child.find_all("li")]
            parts.extend(["- " + item for item in items if item])
        elif child.name == "div":
            # Nested divs (e.g. pilar-question)
            if "pilar-question" in (child.get("class") or []):
                footer = clean(child)
            # Recurse into structural divs (recap-keep, etc.)
            elif "recap-keep-block" not in (child.get("class") or []):
                for p in child.find_all("p"):
                    txt = clean(p)
                    if txt:
                        parts.append(txt)
                for ul in child.find_all("ul"):
                    items = [clean(li) for li in ul.find_all("li")]
                    parts.extend(["- " + item for item in items if item])

    return "\n".join(parts), footer


def card(slide, left, top, w, h, card_title="", body="", footer="",
         accent=ACCENT, title_sz=12, body_sz=10):
    """Simple rounded card with title + body + optional footer."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.04

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.08)
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)

    p = tf.paragraphs[0]
    if card_title:
        run(p, card_title, sz=title_sz, color=accent, bold=True)

    if body:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(5)
        run(p2, body, sz=body_sz, color=TEXT_BODY)

    if footer:
        p3 = tf.add_paragraph()
        p3.space_before = Pt(5)
        run(p3, footer, sz=9, color=TEXT_MUTED, italic=True)


def styled_table(slide, top, headers, rows, w=None):
    if not rows:
        return top
    w = w or CW
    nc = len(headers) if headers else len(rows[0])
    nr = len(rows) + (1 if headers else 0)
    rh = Inches(0.32)
    th = rh * nr

    tbl = slide.shapes.add_table(nr, nc, ML, top, w, th).table

    if headers:
        for j, h in enumerate(headers):
            c = tbl.cell(0, j)
            c.text = h
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(9)
                    r.font.color.rgb = ACCENT
                    r.font.bold = True
                    r.font.name = "Consolas"
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)

    off = 1 if headers else 0
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            if j >= nc:
                break
            c = tbl.cell(i + off, j)
            c.text = str(val)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(9)
                    r.font.color.rgb = TEXT_BODY
                    r.font.name = "Calibri"
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 == 0 else RGBColor(0xF8, 0xF8, 0xF8)

    return top + th + Inches(0.08)


# ===========================================================================
# Slide builders
# ===========================================================================

def build_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, 1)
    t = tb(s, ML, Inches(2.2), CW, Inches(0.35))
    p = t.text_frame.paragraphs[0]
    run(p, "— ESTRATEGIA DE PRODUCTO & EXPERIENCIA", sz=12, color=ACCENT, font="Consolas")
    t2 = tb(s, ML, Inches(2.8), CW, Inches(2.0))
    t2.text_frame.word_wrap = True
    p2 = t2.text_frame.paragraphs[0]
    run(p2, "TIMINING CORE:", sz=36, color=TEXT_DARK, bold=True)
    p3 = t2.text_frame.add_paragraph()
    run(p3, "INTELIGENCIA MULTIMODAL", sz=36, color=TEXT_DARK, bold=True)
    p4 = t2.text_frame.add_paragraph()
    run(p4, "PARA LA MINERÍA", sz=36, color=ACCENT, bold=True)


def build_cards(prs, num, stitle, sub, cards, cols=3):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, num)
    title(s, stitle)
    subtitle(s, sub)

    card_top = Inches(1.35) if sub else Inches(1.1)
    if not cards:
        return

    regular = [c for c in cards if not c.get("is_highlight")]
    highlights = [c for c in cards if c.get("is_highlight")]

    gap = Inches(0.12)
    cw = (CW - gap * (cols - 1)) / cols
    rows_n = (len(regular) + cols - 1) // cols if regular else 0
    avail = SLIDE_H - card_top - Inches(0.3)
    if highlights:
        avail -= Inches(0.55)
    rg = Inches(0.1)
    ch = min((avail - rg * max(0, rows_n - 1)) / max(1, rows_n), Inches(2.8))

    for i, c in enumerate(regular):
        r = i // cols
        col = i % cols
        left = ML + (cw + gap) * col
        top_y = card_top + (ch + rg) * r
        card(s, left, top_y, cw, ch,
             card_title=c.get("title", ""), body=c.get("body", ""),
             footer=c.get("footer", ""), accent=c.get("accent", ACCENT))

    for j, hl in enumerate(highlights):
        hl_top = card_top + (ch + rg) * rows_n + Inches(0.08) + Inches(0.45) * j
        t = tb(s, ML, hl_top, CW, Inches(0.4))
        t.text_frame.word_wrap = True
        p = t.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run(p, hl.get("body", "") or hl.get("title", ""), sz=12, color=ACCENT, font="Consolas")


def build_benchmark(prs, num, stitle, pilar_desc, bench_cards, footer_text=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, num)
    title(s, stitle)
    if pilar_desc:
        subtitle(s, pilar_desc, sz=11, color=TEXT_BODY)

    top_y = Inches(1.4)
    cols = min(len(bench_cards), 3)
    gap = Inches(0.12)
    cw = (CW - gap * (cols - 1)) / cols
    ch = Inches(5.4)

    for i, bc in enumerate(bench_cards[:3]):
        left = ML + (cw + gap) * i
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top_y, cw, ch)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = CARD_BORDER
        shape.line.width = Pt(1)
        shape.adjustments[0] = 0.04

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.15)
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)

        # Name
        p = tf.paragraphs[0]
        run(p, bc.get("name", ""), sz=14, color=TEXT_DARK, bold=True)

        # Industry
        if bc.get("industry"):
            p2 = tf.add_paragraph()
            run(p2, bc["industry"], sz=9, color=TEXT_MUTED)

        # Factor X
        if bc.get("factorx"):
            p3 = tf.add_paragraph()
            p3.space_before = Pt(8)
            run(p3, bc["factorx"], sz=10, color=TEXT_BODY)

        # Image — embed if available, else placeholder text
        if bc.get("image"):
            img_path = f"img/bench/{bc['image']}"
            img_h = Inches(1.8)
            img_top = top_y + Inches(1.6)
            img_left = left + Inches(0.08)
            img_w = cw - Inches(0.16)
            pic = add_image_safe(s, img_path, img_left, img_top, width=img_w)
            if not pic:
                p_img = tf.add_paragraph()
                p_img.space_before = Pt(8)
                run(p_img, f"[img: {bc['image']}]", sz=9, color=TEXT_MUTED, italic=True, font="Consolas")

        # Aplicaciones
        apps = bc.get("aplicaciones", [])
        if apps:
            p4 = tf.add_paragraph()
            p4.space_before = Pt(8)
            run(p4, "Cómo Aplica a CORE:", sz=10, color=ACCENT, bold=True)
            for app in apps[:4]:
                pa = tf.add_paragraph()
                pa.space_before = Pt(2)
                run(pa, f"• {app}", sz=9, color=TEXT_BODY)

    if footer_text:
        t = tb(s, ML, SLIDE_H - Inches(0.45), CW, Inches(0.35))
        t.text_frame.word_wrap = True
        p = t.text_frame.paragraphs[0]
        run(p, footer_text, sz=9, color=TEXT_MUTED)


def build_tables(prs, num, stitle, tables_data):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, num)
    title(s, stitle)
    y = Inches(1.05)
    for td in tables_data:
        sub = td.get("subtitle", "")
        if sub:
            t = tb(s, ML, y, CW, Inches(0.3))
            p = t.text_frame.paragraphs[0]
            run(p, sub, sz=11, color=TEXT_DARK, bold=True)
            y += Inches(0.3)
        y = styled_table(s, y, td.get("headers", []), td.get("rows", []))
        y += Inches(0.05)


def build_rec_table(prs, num, stitle, sub, rec_boxes, table_data):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, num)
    title(s, stitle)
    subtitle(s, sub)

    y = Inches(1.35)
    cols = 2
    gap = Inches(0.1)
    bw = (CW - gap) / 2
    bh = Inches(1.15)
    colors = [ACCENT, COLOR_A, COLOR_A, ACCENT]

    for i, rb in enumerate(rec_boxes[:4]):
        r = i // cols
        c = i % cols
        left = ML + (bw + gap) * c
        top_y = y + (bh + gap) * r
        bc = colors[i] if i < len(colors) else CARD_BORDER
        card(s, left, top_y, bw, bh,
             card_title=rb.get("title", ""), body=rb.get("body", ""),
             accent=bc, title_sz=11, body_sz=10)

    if table_data:
        ty = y + (bh + gap) * 2 + Inches(0.1)
        styled_table(s, ty, table_data.get("headers", []), table_data.get("rows", []))


def build_transition(prs, num, columns):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, num)

    cols_n = len(columns)
    gap = Inches(0.12)
    cw = (CW - gap * (cols_n - 1)) / cols_n
    ch = Inches(5.0)
    top_y = Inches(1.2)
    borders = [CARD_BORDER, ACCENT, TEXT_MUTED]
    title_colors = [TEXT_DARK, ACCENT, TEXT_MUTED]

    for i, c in enumerate(columns[:3]):
        left = ML + (cw + gap) * i
        bc = borders[i] if i < len(borders) else CARD_BORDER
        tc = title_colors[i] if i < len(title_colors) else TEXT_DARK

        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top_y, cw, ch)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = bc
        shape.line.width = Pt(1.5 if i == 1 else 1)
        shape.adjustments[0] = 0.04

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.2)
        tf.margin_left = Inches(0.18)
        tf.margin_right = Inches(0.18)

        tag = c.get("tag", "")
        if tag:
            p = tf.paragraphs[0]
            run(p, tag, sz=10, color=TEXT_MUTED, font="Consolas")

        t = c.get("title", "")
        if t:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(10)
            run(p2, t, sz=18, color=tc, bold=True)

        desc = c.get("desc", "")
        if desc:
            p3 = tf.add_paragraph()
            p3.space_before = Pt(10)
            run(p3, desc, sz=11, color=TEXT_BODY)

        sub = c.get("sub", "")
        if sub:
            p4 = tf.add_paragraph()
            p4.space_before = Pt(8)
            run(p4, sub, sz=10, color=TEXT_MUTED, italic=True)


def build_concentric(prs, num, stitle):
    """Simplified concentric: 3 nested labeled boxes."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, num)
    title(s, stitle)

    cx = SLIDE_W // 2
    cy = Inches(3.8)
    layers = [
        ("VISIÓN ESTRATÉGICA", Inches(10), Inches(4.2), ACCENT,
         "Paz Mental  ·  Dashboard→Copiloto  ·  Valor Percibido"),
        ("DOMINIOS OPERATIVOS", Inches(7.5), Inches(3.0), CARD_BORDER,
         "Respuesta a Crisis  ·  Anticipación Operativa  ·  Asistencia Inteligente  ·  Alineación Estratégica"),
        ("PILARES DE DISEÑO", Inches(5), Inches(1.6), ACCENT,
         "Quiet UI  ·  Clear Path  ·  Time Sight  ·  Omni Sense"),
    ]
    fills = [RGBColor(0xF0, 0xF4, 0xF6), RGBColor(0xE8, 0xEE, 0xF0), RGBColor(0xE0, 0xEA, 0xEE)]

    for i, (label, w, h, border, items) in enumerate(layers):
        shape = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, cx - w // 2, cy - h // 2, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fills[i]
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5 if i == 2 else 1)
        shape.adjustments[0] = 0.06

        # Label top-left
        t = tb(s, cx - w // 2 + Inches(0.25), cy - h // 2 + Inches(0.1), Inches(3), Inches(0.25))
        p = t.text_frame.paragraphs[0]
        run(p, label, sz=9, color=border if border != CARD_BORDER else TEXT_MUTED,
            bold=True, font="Consolas")

        # Items bottom-center
        t2 = tb(s, cx - w // 2, cy + h // 2 - Inches(0.4), w, Inches(0.3))
        p2 = t2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run(p2, items, sz=10, color=TEXT_BODY)

    # DAR label
    dar_top = cy + Inches(2.4)
    t3 = tb(s, ML, dar_top, CW, Inches(0.3))
    p3 = t3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run(p3, "D", sz=12, color=COLOR_D, bold=True, font="Consolas")
    run(p3, " Detección  →  ", sz=10, color=TEXT_MUTED)
    run(p3, "A", sz=12, color=COLOR_A, bold=True, font="Consolas")
    run(p3, " Análisis  →  ", sz=10, color=TEXT_MUTED)
    run(p3, "R", sz=12, color=COLOR_R, bold=True, font="Consolas")
    run(p3, " Recomendación", sz=10, color=TEXT_MUTED)

    t4 = tb(s, ML, dar_top + Inches(0.35), CW, Inches(0.25))
    p4 = t4.text_frame.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    run(p4, "Patrón transversal que conecta todas las capas", sz=9, color=TEXT_MUTED, italic=True)


def build_case_narrative(prs, num, stitle, data):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, num)
    title(s, stitle, sz=18)

    domain = data.get("domain", "")
    context = data.get("context", "")
    perfil = data.get("perfil", {})
    agrupados = data.get("agrupados", [])
    stages = data.get("stages", [])

    # Domain label
    if domain:
        t = tb(s, ML, Inches(0.52), CW, Inches(0.22))
        p = t.text_frame.paragraphs[0]
        run(p, domain, sz=10, color=ACCENT, font="Consolas")

    # Layout: left sidebar (profile) + right (context + DAR stages)
    left_w = Inches(3.2)
    gap = Inches(0.15)
    right_w = CW - left_w - gap
    right_left = ML + left_w + gap
    body_top = Inches(0.9)

    # ── Left sidebar: profile card ──
    left_h = Inches(6.2)
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ML, body_top, left_w, left_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.04
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.12)
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)

    # Roles
    roles = perfil.get("roles", [])
    if roles:
        p = tf.paragraphs[0]
        run(p, "Perfil", sz=10, color=ACCENT, bold=True, font="Consolas")
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        run(p2, " · ".join(roles), sz=9, color=TEXT_DARK, bold=True)

    # Perfil items (Rol, Necesidad, Dolor)
    for item in perfil.get("items", []):
        pi = tf.add_paragraph()
        pi.space_before = Pt(4)
        run(pi, item, sz=9, color=TEXT_BODY)

    # Agrupados
    if agrupados:
        pa = tf.add_paragraph()
        pa.space_before = Pt(10)
        run(pa, "Casos Agrupados", sz=10, color=ACCENT, bold=True, font="Consolas")
        pa2 = tf.add_paragraph()
        pa2.space_before = Pt(4)
        run(pa2, " · ".join(agrupados), sz=9, color=TEXT_BODY)

    # ── Right side: context + 3 DAR stages ──
    # Context
    if context:
        ctx_h = Inches(1.0)
        t = tb(s, right_left, body_top, right_w, ctx_h)
        t.text_frame.word_wrap = True
        p = t.text_frame.paragraphs[0]
        run(p, context, sz=10, color=TEXT_BODY)
        stages_top = body_top + ctx_h + Inches(0.05)
    else:
        stages_top = body_top

    # DAR stages (3 columns)
    cols = min(len(stages), 3)
    sgap = Inches(0.1)
    sw = (right_w - sgap * (cols - 1)) / cols
    sh = SLIDE_H - stages_top - Inches(0.35)
    stage_colors = [COLOR_D, COLOR_A, COLOR_R]
    stage_labels = ["Detección", "Análisis", "Recomendación"]

    for i, stage in enumerate(stages[:3]):
        left = right_left + (sw + sgap) * i
        color = stage_colors[i]

        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, stages_top, sw, sh)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = color
        shape.line.width = Pt(1.5)
        shape.adjustments[0] = 0.04

        stf = shape.text_frame
        stf.word_wrap = True
        stf.margin_top = Inches(0.08)
        stf.margin_left = Inches(0.08)
        stf.margin_right = Inches(0.08)

        p = stf.paragraphs[0]
        run(p, f"E{i+1}: {stage_labels[i]}", sz=11, color=color, bold=True, font="Consolas")

        desc = stage.get("description", "")
        if desc:
            p2 = stf.add_paragraph()
            p2.space_before = Pt(5)
            run(p2, desc, sz=9, color=TEXT_BODY)

        pilar = stage.get("pilar", "")
        if pilar:
            p3 = stf.add_paragraph()
            p3.space_before = Pt(5)
            run(p3, f"Pilar: {pilar}", sz=9, color=ACCENT, italic=True)

        ref_name = stage.get("ref_name", "")
        if ref_name:
            p4 = stf.add_paragraph()
            p4.space_before = Pt(6)
            run(p4, ref_name, sz=9, color=TEXT_DARK, bold=True)
            ref_ind = stage.get("ref_industry", "")
            if ref_ind:
                p5 = stf.add_paragraph()
                run(p5, ref_ind, sz=8, color=TEXT_MUTED)
            ref_fx = stage.get("ref_factorx", "")
            if ref_fx:
                p6 = stf.add_paragraph()
                p6.space_before = Pt(2)
                run(p6, ref_fx, sz=8, color=TEXT_MUTED, italic=True)

        for app in stage.get("aplicaciones", [])[:3]:
            pa = stf.add_paragraph()
            pa.space_before = Pt(2)
            run(pa, f"• {app}", sz=8, color=TEXT_BODY)


def build_case_viz(prs, num, stitle, data):
    """VIZ slides: DAR indicator + moment card + pilars + refs (no image)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, num)
    title(s, stitle, sz=16)

    active = data.get("active", "D")

    # DAR bar
    dar_top = Inches(0.9)
    t = tb(s, ML, dar_top, CW, Inches(0.3))
    p = t.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for letter, label, color in [("D", "Detección", COLOR_D),
                                  ("A", "Análisis", COLOR_A),
                                  ("R", "Recomendación", COLOR_R)]:
        is_active = (letter == active)
        c = color if is_active else TEXT_MUTED
        run(p, f"  {letter} {label}  ", sz=11, color=c,
            bold=is_active, font="Consolas")

    # Two columns: left = placeholder, right = info cards
    col_gap = Inches(0.2)
    left_w = CW * 0.55
    right_w = CW - left_w - col_gap
    right_left = ML + left_w + col_gap
    info_top = Inches(1.5)

    # Left: viz image or placeholder
    viz_img = data.get("viz_image", "")
    pic = add_image_safe(s, viz_img, ML, info_top, width=left_w) if viz_img else None
    if not pic:
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ML, info_top, left_w, Inches(5.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
        shape.line.color.rgb = CARD_BORDER
        shape.line.width = Pt(1)
        shape.adjustments[0] = 0.03
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = tf.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        run(pp, f"[Visualización: {viz_img}]" if viz_img else "[Visualización]",
            sz=11, color=TEXT_MUTED, font="Consolas")

    # Right: moment card
    mt = data.get("moment_title", "")
    md = data.get("moment_desc", "")
    card(s, right_left, info_top, right_w, Inches(1.8),
         card_title=mt, body=md, title_sz=11, body_sz=10)

    # Right: pilars card
    pt = data.get("pilars_text", "")
    card(s, right_left, info_top + Inches(1.9), right_w, Inches(1.2),
         card_title="Pilares", body=pt, accent=ACCENT, title_sz=10, body_sz=9)

    # Right: refs
    refs = data.get("refs", [])
    ref_top = info_top + Inches(3.2)
    t2 = tb(s, right_left, ref_top, right_w, Inches(0.25))
    p2 = t2.text_frame.paragraphs[0]
    run(p2, "REFERENTES POR ETAPA", sz=9, color=ACCENT, font="Consolas")

    ref_colors = [COLOR_D, COLOR_A, COLOR_R]
    ref_letters = ["D", "A", "R"]
    for i, ref in enumerate(refs[:3]):
        ry = ref_top + Inches(0.3) + Inches(0.35) * i
        card(s, right_left, ry, right_w, Inches(0.3),
             card_title=f"{ref_letters[i]}  {ref}", accent=ref_colors[i],
             title_sz=9, body_sz=9)


def build_text(prs, num, stitle, paragraphs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, num)
    title(s, stitle)
    t = tb(s, ML, Inches(1.1), CW, Inches(5.5))
    t.text_frame.word_wrap = True
    for i, txt in enumerate(paragraphs):
        p = t.text_frame.paragraphs[0] if i == 0 else t.text_frame.add_paragraph()
        if i > 0:
            p.space_before = Pt(8)
        run(p, txt, sz=11, color=TEXT_BODY)


def build_idemax(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, IDEMAX_RED)
    logo_path = BASE_DIR / "logo_idemax.png"
    if logo_path.exists():
        img_w = Inches(3.5)
        s.shapes.add_picture(str(logo_path), (SLIDE_W - img_w) // 2, Inches(3.0), img_w)
    else:
        t = tb(s, Inches(3), Inches(3), Inches(7), Inches(1.5))
        p = t.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run(p, "IDEMAX", sz=48, color=WHITE, bold=True)


# ===========================================================================
# HTML Parser (reused from original, simplified)
# ===========================================================================

def parse_bench_card(bc_el):
    name_el = bc_el.find("div", class_="bench-name")
    industry_el = bc_el.find("div", class_="bench-industry")
    factorx_el = bc_el.find("div", class_="bench-factorx")
    img_el = bc_el.find("div", class_="bench-ui-placeholder")
    apps_el = bc_el.find("ul", class_="bench-steal")

    img_file = ""
    if img_el:
        img = img_el.find("img")
        if img and img.get("src", "").startswith("img/bench/"):
            img_file = img["src"].split("/")[-1]

    aplicaciones = []
    if apps_el:
        for li in apps_el.find_all("li"):
            aplicaciones.append(clean(li))

    return {
        "name": clean(name_el), "industry": clean(industry_el),
        "factorx": clean(factorx_el), "image": img_file,
        "aplicaciones": aplicaciones,
    }


def parse_table(table_el):
    headers = []
    rows = []
    thead = table_el.find("thead")
    if thead:
        headers = [clean(th) for th in thead.find_all("th")]
    tbody = table_el.find("tbody")
    tr_source = tbody.find_all("tr") if tbody else table_el.find_all("tr")
    for tr in tr_source:
        cells = [clean(td) for td in tr.find_all(["td", "th"])]
        if cells and cells != headers:
            rows.append(cells)
    return headers, rows


def parse_slides(html_path):
    with open(html_path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f, "html.parser")

    slides = soup.find_all("section", class_="slide")
    parsed = []

    for i, section in enumerate(slides):
        sid = section.get("id", f"slide-{i+1}")
        num = i + 1
        title_el = section.find("h2", class_="slide-title")
        stitle = clean(title_el) if title_el else ""
        if not stitle or stitle == "—":
            h1 = section.find("h1")
            stitle = clean(h1) if h1 else f"Slide {num}"

        cards_el = section.find_all("div", class_="card")
        bench_el = section.find_all("div", class_="bench-card")
        case_tl = section.find("div", class_="case-timeline")
        has_viz = any(img.get("src", "").startswith("viz-") for img in section.find_all("img"))
        trans_cols = section.find_all("div", class_="transition-col")
        layer_rings = section.find_all("div", class_="layer-ring")
        rec_boxes = section.find_all("div", class_="rec-box")
        tables = section.find_all("table")
        horizon = section.find_all("div", class_=lambda c: c and "horizon-ladder-step" in c)
        cg = section.find("div", class_="content-grid")

        info = {"num": num, "id": sid, "title": stitle, "type": "text", "data": {}}

        # Idemax
        if "slide-idemax" in (section.get("class") or []):
            info["type"] = "idemax"

        # Cover
        elif num == 1:
            info["type"] = "cover"

        # Case VIZ
        elif has_viz:
            info["type"] = "case_viz"
            viz_img = ""
            for img in section.find_all("img"):
                if img.get("src", "").startswith("viz-"):
                    viz_img = img["src"]
                    break
            active = "D"
            for span in section.find_all("span"):
                style = span.get("style", "")
                text = span.get_text()
                if "Recomendación" in text and "var(--accent-cyan)" in style and "0.3" not in style:
                    active = "R"; break
                elif "Análisis" in text and "var(--color-a)" in style and "0.3" not in style:
                    active = "A"; break
                elif "Detección" in text and "var(--color-d)" in style and "0.3" not in style:
                    active = "D"; break

            sidebar = section.find_all("div", class_="card")
            mt = md = pt = ""
            if len(sidebar) >= 1:
                h3 = sidebar[0].find("h3")
                mt = clean(h3) if h3 else ""
                p_el = sidebar[0].find("p")
                md = clean(p_el) if p_el else ""
            if len(sidebar) >= 2:
                h3_2 = sidebar[1].find("h3")
                p_2 = sidebar[1].find("p")
                pt = (clean(h3_2) + ": " + clean(p_2)) if h3_2 else clean(p_2)

            refs = []
            for rs in section.find_all("span", style=lambda s: s and "font-size: 12px" in s and "color: #ccc" in s):
                refs.append(clean(rs))

            info["data"] = {"viz_image": viz_img, "active": active,
                            "moment_title": mt, "moment_desc": md,
                            "pilars_text": pt, "refs": refs}

        # Case narrative
        elif case_tl:
            info["type"] = "case_narrative"

            # Domain
            domain_el = section.find("div", class_="case-dominio")

            # Context (in case-contexto-bajada, not case-ctx)
            ctx_el = section.find("div", class_="case-contexto-bajada")
            if not ctx_el:
                ctx_el = section.find("p", class_="case-ctx")

            # Profile: roles, perfil items, agrupados
            roles = [clean(r) for r in section.find_all("span", class_="cluster-role")]
            perfil_items = [clean(pi) for pi in section.find_all("div", class_="case-perfil-item")]
            agrupados_el = section.find("div", class_="case-agrupados")
            agrupados = [clean(s) for s in agrupados_el.find_all("span")] if agrupados_el else []

            # DAR stages
            stages = []
            for se in case_tl.find_all("div", class_="case-stage"):
                stages.append({
                    "description": clean(se.find("p", class_="case-descripcion")),
                    "pilar": clean(se.find("span", class_="case-pilar-name")),
                    "ref_name": clean(se.find("div", class_="case-ref-name")),
                    "ref_industry": clean(se.find("div", class_="case-ref-industry")),
                    "ref_factorx": clean(se.find("div", class_="case-ref-factorx")),
                    "aplicaciones": [clean(li) for li in se.find_all("li")]
                })
            info["data"] = {
                "domain": clean(domain_el),
                "context": clean(ctx_el),
                "perfil": {"roles": roles, "items": perfil_items},
                "agrupados": agrupados,
                "stages": stages,
            }

        # Transition
        elif trans_cols:
            info["type"] = "transition"
            columns = []
            for tc in trans_cols:
                columns.append({
                    "tag": clean(tc.find("div", class_="transition-tag")),
                    "title": clean(tc.find("h3", class_="transition-title")),
                    "desc": clean(tc.find("p", class_="transition-desc")),
                    "sub": clean(tc.find("p", class_="transition-sub")),
                })
            info["data"] = {"columns": columns}

        # Benchmark
        elif bench_el:
            info["type"] = "benchmark"
            bench_data = [parse_bench_card(bc) for bc in bench_el]
            pilar_h = section.find("div", class_="bench-pilar-header")
            pd = ""
            if pilar_h:
                n = clean(pilar_h.find("div", class_="bench-pilar-name"))
                tl = clean(pilar_h.find("div", class_="bench-pilar-tagline"))
                pd = f"{n}. {tl}" if tl else n
            footer = ""
            for p in section.find_all("p"):
                if "También" in clean(p):
                    footer = clean(p)
                    break
            info["data"] = {"cards": bench_data, "pilar_desc": pd, "footer": footer}

        # Two tables
        elif len(tables) >= 2:
            info["type"] = "two_tables"
            td = []
            for tbl in tables:
                h, r = parse_table(tbl)
                prev = tbl.find_previous_sibling("h3")
                if not prev:
                    parent = tbl.parent
                    if parent:
                        prev = parent.find("h3")
                td.append({"subtitle": clean(prev) if prev else "", "headers": h, "rows": r})
            info["data"] = {"tables": td}

        # Rec-boxes + table
        elif rec_boxes and tables:
            info["type"] = "rec_table"
            boxes = []
            for rb in rec_boxes:
                ts = rb.find("span", class_="rec-title")
                rt = clean(ts) if ts else ""
                rb_body = clean(rb)
                if rt:
                    rb_body = rb_body.replace(rt, "", 1).strip()
                boxes.append({"title": rt, "body": rb_body})
            h, r = parse_table(tables[0])
            sub = ""
            if cg:
                for div in cg.find_all("div", class_="col-12", recursive=False):
                    fp = div.find("p")
                    if fp:
                        sub = clean(fp)
                        break
            info["data"] = {"subtitle": sub, "rec_boxes": boxes,
                            "table": {"headers": h, "rows": r}}

        # Concentric
        elif layer_rings:
            info["type"] = "concentric"

        # Cards grid
        elif cards_el:
            info["type"] = "cards"
            sub = ""
            if cg:
                for el in cg.children:
                    if not hasattr(el, 'find') or not hasattr(el, 'get'):
                        continue
                    if el.find("p") and "card" not in (el.get("class") or []):
                        p_el = el.find("p")
                        if p_el and not p_el.find_parent("div", class_="card"):
                            candidate = clean(p_el)
                            if candidate and len(candidate) > 20:
                                sub = candidate
                                break

            col_classes = set()
            for c in cards_el:
                for cls in c.get("class", []):
                    if cls.startswith("col-"):
                        col_classes.add(cls)
            cols = 3
            if "col-3" in col_classes:
                cols = 4
            elif "col-6" in col_classes and "col-4" not in col_classes:
                cols = 2

            card_data = []

            # Check for recap-keep-block (slide 02) — inject as first card
            keep_block = section.find("div", class_="recap-keep-block")
            if keep_block:
                keep_items = keep_block.find_all("div", class_="recap-keep-item")
                keep_parts = []
                for ki in keep_items:
                    strong = ki.find("strong")
                    span = ki.find("span")
                    s_text = clean(strong) if strong else ""
                    d_text = clean(span) if span else ""
                    if s_text and d_text:
                        keep_parts.append(f"{s_text}: {d_text}")
                    elif s_text:
                        keep_parts.append(s_text)
                keep_header_el = keep_block.find("div", class_="recap-keep-header")
                keep_title = clean(keep_header_el) if keep_header_el else "KEEP"
                card_data.append({
                    "title": keep_title,
                    "body": " | ".join(keep_parts),
                    "footer": "", "is_highlight": False,
                })
                cols = max(cols, 1)  # full width for keep block row

            for c in cards_el:
                ct = ""
                h3 = c.find("h3")
                h4 = c.find("h4")
                if h3:
                    ct = clean(h3)
                elif h4:
                    ct = clean(h4)
                body_text, cf = extract_body(c)
                is_hl = False
                if "col-12" in c.get("class", []):
                    style = c.get("style", "")
                    if "dashed" in style or "transparent" in style:
                        is_hl = True
                card_data.append({"title": ct, "body": body_text,
                                  "footer": cf, "is_highlight": is_hl})

            # Capture orphan rec-boxes (not inside cards) as highlight cards
            if rec_boxes:
                for rb in rec_boxes:
                    if not rb.find_parent("div", class_="card"):
                        ts = rb.find("span", class_="rec-title")
                        rt = clean(ts) if ts else ""
                        rb_body = clean(rb)
                        if rt:
                            rb_body = rb_body.replace(rt, "", 1).strip()
                        card_data.append({"title": rt, "body": rb_body,
                                          "footer": "", "is_highlight": True})

            info["data"] = {"subtitle": sub, "cols": cols, "cards": card_data}

        # Horizon steps
        elif horizon:
            info["type"] = "text"
            paragraphs = []
            for hs in horizon:
                label = clean(hs.find("div", class_="horizon-ladder-label"))
                text = clean(hs.find("div", class_="horizon-ladder-text"))
                paragraphs.append(f"{label}: {text}")
            info["data"] = {"paragraphs": paragraphs}

        # Single table
        elif tables:
            info["type"] = "single_table"
            h, r = parse_table(tables[0])
            info["data"] = {"headers": h, "rows": r}

        # Fallback
        else:
            paragraphs = []
            if cg:
                for p in cg.find_all("p"):
                    txt = clean(p)
                    if txt:
                        paragraphs.append(txt)
            if not paragraphs:
                all_t = clean(section)
                if all_t:
                    paragraphs.append(all_t)
            info["data"] = {"paragraphs": paragraphs if paragraphs else [""]}

        parsed.append(info)

    return parsed


# ===========================================================================
# Main
# ===========================================================================

def generate(slides_data, output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for sd in slides_data:
        t = sd["type"]
        n = sd["num"]
        st = sd["title"]
        d = sd["data"]

        if t == "idemax":
            build_idemax(prs)
        elif t == "cover":
            build_cover(prs)
        elif t == "cards":
            build_cards(prs, n, st, d.get("subtitle", ""),
                        d.get("cards", []), cols=d.get("cols", 3))
        elif t == "benchmark":
            build_benchmark(prs, n, st, d.get("pilar_desc", ""),
                            d.get("cards", []), d.get("footer", ""))
        elif t == "two_tables":
            build_tables(prs, n, st, d.get("tables", []))
        elif t == "rec_table":
            build_rec_table(prs, n, st, d.get("subtitle", ""),
                            d.get("rec_boxes", []), d.get("table", {}))
        elif t == "transition":
            build_transition(prs, n, d.get("columns", []))
        elif t == "concentric":
            build_concentric(prs, n, st)
        elif t == "case_narrative":
            build_case_narrative(prs, n, st, d)
        elif t == "case_viz":
            build_case_viz(prs, n, st, d)
        elif t == "single_table":
            s = prs.slides.add_slide(prs.slide_layouts[6])
            bg(s)
            header(s, n)
            title(s, st)
            styled_table(s, Inches(1.1), d.get("headers", []), d.get("rows", []))
        elif t == "text":
            build_text(prs, n, st, d.get("paragraphs", [""]))
        else:
            build_text(prs, n, st, [f"[{t}]"])

    prs.save(str(output_path))
    return len(prs.slides)


if __name__ == "__main__":
    print(f"Parsing {HTML_FILE}...")
    data = parse_slides(HTML_FILE)
    print(f"Found {len(data)} slides\n")

    types = {}
    for sd in data:
        t = sd["type"]
        types[t] = types.get(t, 0) + 1
        print(f"  S{sd['num']:02d} [{t:15s}] {sd['title'][:55]}")

    print(f"\nTypes: {dict(sorted(types.items(), key=lambda x: -x[1]))}")
    print(f"\nGenerating PPTX (light mode)...")
    n = generate(data, OUTPUT_FILE)
    print(f"Done! {n} slides → {OUTPUT_FILE}")
    print(f"Size: {OUTPUT_FILE.stat().st_size / 1024:.0f} KB")
